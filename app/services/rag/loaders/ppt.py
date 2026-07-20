"""PPTX slide, notes, and table extraction."""
from pathlib import Path
import shutil
import subprocess
import tempfile
from loguru import logger


def extract_ppt_text(file_path: str | Path) -> list[dict]:
    path = Path(file_path)
    if path.suffix.lower() == ".ppt":
        converter = shutil.which("soffice") or shutil.which("libreoffice")
        if not converter:
            raise ValueError("Legacy .ppt ingestion requires LibreOffice (soffice) on PATH")
        with tempfile.TemporaryDirectory(prefix="follei-ppt-") as temp_dir:
            completed = subprocess.run([converter, "--headless", "--convert-to", "pptx", "--outdir", temp_dir, str(path)], capture_output=True, text=True, timeout=90, check=False)
            converted = Path(temp_dir) / f"{path.stem}.pptx"
            if completed.returncode or not converted.is_file():
                raise ValueError(f"Legacy .ppt conversion failed: {completed.stderr.strip()[:500]}")
            return extract_ppt_text(converted)
    from pptx import Presentation
    deck = Presentation(str(path))
    pages = []
    for number, slide in enumerate(deck.slides, start=1):
        parts, heading = [], None
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                value = shape.text.strip(); heading = heading or value.splitlines()[0]; parts.append(value)
            if getattr(shape, "has_table", False):
                parts.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes:\n{notes}")
        if parts:
            pages.append({"page": number, "heading": heading, "text": "\n".join(parts)})
    logger.info(f"Extracted {len(pages)} slides from {path.name}")
    return pages
