from pathlib import Path
import sys
from types import SimpleNamespace
from app.services.rag.loaders.email import extract_email_text
from app.services.rag.loaders.ppt import extract_ppt_text
from app.services.rag.loaders.spreadsheet import extract_csv_text, extract_xlsx_text


def test_csv_loader_preserves_header_and_rows(tmp_path):
    path = tmp_path / "pricing.csv"; path.write_text("Plan,Price\nEnterprise,999\n", encoding="utf-8")
    assert "Plan | Price" in extract_csv_text(path)[0]["text"]


def test_xlsx_loader_preserves_sheet_heading(tmp_path):
    from openpyxl import Workbook
    path = tmp_path / "pricing.xlsx"; book = Workbook(); sheet = book.active; sheet.title = "Plans"; sheet.append(["Plan", "Price"]); sheet.append(["Enterprise", 999]); book.save(path)
    assert extract_xlsx_text(path)[0]["heading"] == "Plans"


def test_pptx_loader_extracts_slide_text(tmp_path):
    from pptx import Presentation
    path = tmp_path / "deck.pptx"; deck = Presentation(); slide = deck.slides.add_slide(deck.slide_layouts[1]); slide.shapes.title.text = "Enterprise Plan"; slide.placeholders[1].text = "USD 999 per month"; slide.notes_slide.notes_text_frame.text = "Mention the 45-day refund window."; deck.save(path)
    extracted = extract_ppt_text(path)[0]["text"]
    assert "USD 999" in extracted
    assert "Mention the 45-day refund window." in extracted


def test_eml_loader_extracts_bounded_text_attachment(tmp_path):
    from email.message import EmailMessage
    path = tmp_path / "mail.eml"; message = EmailMessage(); message["Subject"] = "Refund"; message.set_content("Refund policy attached"); message.add_attachment(b"window,45 days\n", maintype="text", subtype="csv", filename="policy.csv"); path.write_bytes(message.as_bytes())
    assert "45 days" in extract_email_text(path)[0]["text"]


def test_msg_loader_extracts_only_bounded_text_attachments(tmp_path, monkeypatch):
    path = tmp_path / "mail.msg"
    path.write_bytes(b"fake-msg-container")
    closed = {}

    class FakeMessage:
        subject = "Pricing"
        sender = "sales@example.com"
        body = "See the attached terms."
        attachments = [
            SimpleNamespace(longFilename="plans.csv", shortFilename=None, data=b"plan,price\nEnterprise,999\n"),
            SimpleNamespace(longFilename="oversized.txt", shortFilename=None, data=b"x" * (2 * 1024 * 1024 + 1)),
            SimpleNamespace(longFilename="unsafe.exe", shortFilename=None, data=b"do not extract"),
        ]

        def __init__(self, _path):
            pass

        def close(self):
            closed["value"] = True

    monkeypatch.setitem(sys.modules, "extract_msg", SimpleNamespace(Message=FakeMessage))

    text = extract_email_text(path)[0]["text"]

    assert "Enterprise,999" in text
    assert "oversized" not in text
    assert "do not extract" not in text
    assert closed["value"] is True


def test_scanned_pdf_ocr_passes_a_decoded_image_to_tesseract(tmp_path, monkeypatch):
    import sys
    from PIL import Image
    from app.services.rag.loaders import pdf as pdf_loader

    path = tmp_path / "scan.pdf"
    Image.new("RGB", (300, 120), "white").save(path, "PDF")
    seen = {}

    def fake_ocr(image):
        seen["is_image"] = isinstance(image, Image.Image)
        return "The refund window is 45 days."

    monkeypatch.setattr(pdf_loader.shutil, "which", lambda _name: sys.executable)
    monkeypatch.setattr(pdf_loader.pytesseract, "image_to_string", fake_ocr)

    pages = pdf_loader.extract_pdf_text(path)

    assert seen["is_image"] is True
    assert pages[0]["text"] == "The refund window is 45 days."
